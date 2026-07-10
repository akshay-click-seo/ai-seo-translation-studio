"""
services/pptx_service.py — PowerPoint (.pptx) translation service.

Preserves:
  - Slide layout and design
  - Font styles (bold, italic, size, colour)
  - Text box positions and sizes
  - Slide notes (speaker notes)
  - Tables on slides
  - Grouped shapes

Translates:
  - Text frames in all shapes
  - Table cell text
  - Speaker notes
"""

import io
import os
import time
from typing import Optional, Callable

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from logger import setup_logger, translation_logger
from utils import is_translatable

logger = setup_logger(__name__)


class PptxService:
    """Translates .pptx presentations while preserving visual design."""

    def translate(
        self,
        file_path: str,
        translator,
        target_lang: str = "es-la",
        mode: str = "Technical",
        translate_notes: bool = True,
        use_tm: bool = True,
        on_progress: Optional[Callable[[int, int], None]] = None,
    ) -> bytes:
        """
        Translate all text in a PowerPoint presentation.

        Args:
            file_path:        Path to source .pptx file.
            translator:       Translator instance.
            target_lang:      Target language code.
            mode:             Translation mode.
            translate_notes:  Whether to translate speaker notes.
            use_tm:           Use Translation Memory.
            on_progress:      Progress callback(current_slide, total_slides).

        Returns:
            Translated presentation as raw bytes.
        """
        t0 = time.perf_counter()
        prs = Presentation(file_path)
        total = len(prs.slides)
        segment_count = 0

        for slide_idx, slide in enumerate(prs.slides):
            if on_progress:
                on_progress(slide_idx, total)

            # ── Shapes ────────────────────────────────────────────────────────
            for shape in slide.shapes:
                count = self._translate_shape(
                    shape, translator, target_lang, mode, use_tm
                )
                segment_count += count

            # ── Speaker notes ─────────────────────────────────────────────────
            if translate_notes and slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                for para in notes_frame.paragraphs:
                    self._translate_paragraph(para, translator, target_lang, mode, use_tm)
                    segment_count += 1

        # Serialize
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        result = buf.read()

        elapsed = (time.perf_counter() - t0) * 1000
        translation_logger.log_file_translation(
            file_name=os.path.basename(file_path),
            file_type="pptx",
            mode=mode,
            target_lang=target_lang,
            segments=segment_count,
            processing_time_ms=elapsed,
        )
        logger.info("PPTX translated: %d segments in %.0f ms", segment_count, elapsed)
        return result

    # ── Private helpers ────────────────────────────────────────────────────────

    def _translate_shape(
        self,
        shape,
        translator,
        target_lang: str,
        mode: str,
        use_tm: bool,
    ) -> int:
        """
        Translate text in a single shape (recursing into groups/tables).

        Returns:
            Number of text segments translated.
        """
        count = 0

        # Group shapes — recurse
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for child in shape.shapes:
                count += self._translate_shape(child, translator, target_lang, mode, use_tm)
            return count

        # Table
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        self._translate_paragraph(para, translator, target_lang, mode, use_tm)
                        count += 1
            return count

        # Text frame
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                self._translate_paragraph(para, translator, target_lang, mode, use_tm)
                count += 1

        return count

    def _translate_paragraph(
        self,
        para,
        translator,
        target_lang: str,
        mode: str,
        use_tm: bool,
    ) -> None:
        """
        Translate a single paragraph in-place, preserving run-level formatting.

        Args:
            para:        pptx Paragraph object.
            translator:  Translator instance.
            target_lang: Target language code.
            mode:        Translation mode.
            use_tm:      Use Translation Memory.
        """
        full_text = "".join(run.text for run in para.runs)
        if not full_text.strip() or not is_translatable(full_text):
            return

        try:
            result = translator.translate(
                full_text.strip(), target_lang=target_lang, mode=mode, use_tm=use_tm
            )
            translated = result["translation"]
        except Exception as exc:
            logger.error("PPTX paragraph error: %s", exc)
            return

        if not para.runs:
            return

        if len(para.runs) == 1:
            para.runs[0].text = translated
        else:
            # Preserve first-run formatting; clear other runs
            para.runs[0].text = translated
            for run in para.runs[1:]:
                run.text = ""

    # ── Utilities ──────────────────────────────────────────────────────────────

    @staticmethod
    def extract_text(file_path: str) -> str:
        """
        Extract all text content from a .pptx file.

        Args:
            file_path: Path to the .pptx file.

        Returns:
            All slide text joined by newlines.
        """
        prs = Presentation(file_path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(r.text for r in para.runs).strip()
                        if text:
                            lines.append(text)
        return "\n".join(lines)

    @staticmethod
    def get_stats(file_path: str) -> dict:
        """
        Return slide and text statistics for a .pptx file.

        Returns:
            Dict with keys: slides, shapes, words.
        """
        prs = Presentation(file_path)
        all_text = PptxService.extract_text(file_path)
        return {
            "slides": len(prs.slides),
            "shapes": sum(len(slide.shapes) for slide in prs.slides),
            "words":  len(all_text.split()),
        }


# ── Module-level singleton ─────────────────────────────────────────────────────
pptx_service = PptxService()
